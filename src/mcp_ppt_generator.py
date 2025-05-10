import os
import json
import logging
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import base64
from io import BytesIO
from PIL import Image

# 配置日誌
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("mcp_ppt_generator.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class MCPPPTGenerator:
    """使用MCP協議自動生成PPT的模組"""
    
    def __init__(self, api_key=None, model="claude-3-7-sonnet-20250219"):
        """初始化MCP PPT生成器
        
        Args:
            api_key: Anthropic API密鑰
            model: 使用的Claude模型
        """
        # 獲取API密鑰
        self.api_key = api_key or os.getenv("ANTHROPIC_API_KEY")
        if not self.api_key:
            raise ValueError("需要Anthropic API密鑰")
            
        # 初始化Claude客戶端
        self.client = anthropic.Anthropic(api_key=self.api_key)
        self.model = model
    
    def generate_ppt_from_translation(self, translation_data, output_path):
        """從翻譯數據生成PPT
        
        Args:
            translation_data: 翻譯後的數據
            output_path: PPT輸出路徑
            
        Returns:
            PPT文件路徑
        """
        logger.info(f"從翻譯數據生成PPT：{output_path}")
        
        # 準備翻譯內容
        content = self._prepare_content_from_translation(translation_data)
        
        # 使用MCP協議生成PPT結構
        ppt_structure = self._generate_ppt_structure(content)
        
        # 創建PPT
        self._create_ppt_from_structure(ppt_structure, output_path)
        
        logger.info(f"PPT生成完成：{output_path}")
        return output_path
    
    def _prepare_content_from_translation(self, translation_data):
        """從翻譯數據準備內容
        
        Args:
            translation_data: 翻譯後的數據
            
        Returns:
            準備好的內容字典
        """
        content = {
            "title": translation_data.get("filename", "Untitled Document"),
            "sections": [],
            "images": []
        }
        
        # 提取中文翻譯內容
        for page in translation_data.get("text_data", []):
            for block in page.get("blocks", []):
                if "content_translated" in block and block.get("type") == "text":
                    # 檢查是否是標題（簡單啟發式方法）
                    text = block["content_translated"]
                    if len(text) < 100 and not text.endswith('.'):
                        # 可能是標題
                        content["sections"].append({
                            "title": text,
                            "content": []
                        })
                    else:
                        # 正文內容
                        if not content["sections"]:
                            content["sections"].append({
                                "title": "主要內容",
                                "content": []
                            })
                        content["sections"][-1]["content"].append(text)
        
        # 提取圖像（如果有）
        for img_info in translation_data.get("images", []):
            if "text_in_image_translated" in img_info:
                content["images"].append({
                    "caption": img_info.get("text_in_image_translated", ""),
                    "page": img_info.get("page_num", 0)
                })
        
        return content
    
    def _generate_ppt_structure(self, content):
        """使用MCP協議生成PPT結構
        
        Args:
            content: 準備好的內容
            
        Returns:
            PPT結構
        """
        # 構建MCP提示
        prompt = self._create_mcp_prompt(content)
        
        try:
            # 調用Claude API
            response = self.client.messages.create(
                model=self.model,
                max_tokens=4000,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            
            # 提取JSON結構
            ppt_structure = self._extract_json_from_response(response.content[0].text)
            
            # 驗證結構
            if not ppt_structure or not isinstance(ppt_structure, dict) or "slides" not in ppt_structure:
                logger.error("生成的PPT結構無效")
                # 使用備用結構
                ppt_structure = self._create_fallback_structure(content)
            
            return ppt_structure
            
        except Exception as e:
            logger.error(f"生成PPT結構時出錯: {str(e)}")
            # 使用備用結構
            return self._create_fallback_structure(content)
    
    def _create_mcp_prompt(self, content):
        """創建MCP提示
        
        Args:
            content: 準備好的內容
            
        Returns:
            MCP提示
        """
        prompt = f"""
您是一個專業的演示文稿設計專家。請使用以下內容創建一個精簡且專業的學術簡報，使用MCP（Multi-modal Conversational Protocol）協議。

文檔標題：{content["title"]}

文檔內容：
"""
        
        # 添加章節內容
        for section in content["sections"]:
            prompt += f"\n## {section['title']}\n"
            for paragraph in section["content"]:
                prompt += f"{paragraph}\n"
        
        prompt += """
請創建一個包含以下元素的簡報結構：
1. 封面幻燈片（包含標題）
2. 目錄幻燈片
3. 按章節組織的內容幻燈片
4. 總結幻燈片

對於每張幻燈片，請提供：
- 幻燈片標題
- 幻燈片內容（要點列表或段落）
- 幻燈片設計建議（如布局、顏色等）

請以JSON格式提供響應，結構如下：
```json
{
  "title": "簡報標題",
  "theme": "建議的主題顏色和風格",
  "slides": [
    {
      "type": "cover",
      "title": "標題文本",
      "subtitle": "副標題文本",
      "design": "設計說明"
    },
    {
      "type": "toc",
      "title": "目錄",
      "items": ["項目1", "項目2", ...],
      "design": "設計說明"
    },
    {
      "type": "content",
      "title": "幻燈片標題",
      "bullet_points": ["要點1", "要點2", ...],
      "design": "設計說明"
    },
    ...
    {
      "type": "summary",
      "title": "總結",
      "bullet_points": ["要點1", "要點2", ...],
      "design": "設計說明"
    }
  ]
}
```

請確保JSON格式正確，以便能夠自動處理生成PPT。內容應該是濃縮的學術要點，而非完整文章。每張內容幻燈片應限制在3-5個要點，每個要點不超過2行。
"""
        
        return prompt
    
    def _extract_json_from_response(self, response_text):
        """從響應中提取JSON
        
        Args:
            response_text: API響應文本
            
        Returns:
            提取的JSON結構
        """
        # 使用正則表達式提取JSON部分
        json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
        
        if json_match:
            json_str = json_match.group(1)
        else:
            # 如果沒有找到JSON標記，嘗試直接解析整個響應
            json_str = response_text
        
        # 嘗試解析JSON
        try:
            return json.loads(json_str)
        except json.JSONDecodeError as e:
            logger.error(f"JSON解析錯誤：{str(e)}")
            return None
    
    def _create_fallback_structure(self, content):
        """創建備用PPT結構
        
        Args:
            content: 準備好的內容
            
        Returns:
            備用PPT結構
        """
        title = content.get("title", "Untitled Document")
        
        # 創建基本結構
        structure = {
            "title": title,
            "theme": "專業學術風格，藍色主題",
            "slides": [
                {
                    "type": "cover",
                    "title": title,
                    "subtitle": "學術論文簡報",
                    "design": "中央對齊，藍色背景"
                },
                {
                    "type": "toc",
                    "title": "目錄",
                    "items": [section["title"] for section in content["sections"]],
                    "design": "左對齊目錄項目，右側留白"
                }
            ]
        }
        
        # 為每個章節創建內容幻燈片
        for section in content["sections"]:
            # 提取重要句子作為要點
            bullet_points = []
            for paragraph in section["content"]:
                # 簡單提取第一句話
                sentences = paragraph.split('。')
                if sentences:
                    first_sentence = sentences[0] + '。'
                    if first_sentence not in bullet_points and len(first_sentence) > 10:
                        bullet_points.append(first_sentence)
            
            # 限制要點數量
            bullet_points = bullet_points[:5]
            
            if bullet_points:
                structure["slides"].append({
                    "type": "content",
                    "title": section["title"],
                    "bullet_points": bullet_points,
                    "design": "標準內容幻燈片，左對齊"
                })
        
        # 添加總結幻燈片
        structure["slides"].append({
            "type": "summary",
            "title": "總結",
            "bullet_points": ["文檔主要論點和發現"],
            "design": "強調性設計，深色背景，白色文字"
        })
        
        return structure
    
    def _create_ppt_from_structure(self, structure, output_path):
        """從結構創建PPT
        
        Args:
            structure: PPT結構
            output_path: 輸出路徑
        """
        # 創建演示文稿
        prs = Presentation()
        
        # 設置幻燈片尺寸為寬屏16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 處理每張幻燈片
        for slide_info in structure["slides"]:
            slide_type = slide_info.get("type", "")
            
            if slide_type == "cover":
                self._create_cover_slide(prs, slide_info)
            elif slide_type == "toc":
                self._create_toc_slide(prs, slide_info)
            elif slide_type == "content":
                self._create_content_slide(prs, slide_info)
            elif slide_type == "summary":
                self._create_summary_slide(prs, slide_info)
            else:
                # 默認創建內容幻燈片
                self._create_content_slide(prs, slide_info)
        
        # 保存演示文稿
        prs.save(output_path)
    
    def _create_cover_slide(self, prs, slide_info):
        """創建封面幻燈片
        
        Args:
            prs: 演示文稿對象
            slide_info: 幻燈片信息
        """
        # 使用標題幻燈片佈局
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # 設置標題
        title = slide.shapes.title
        title.text = slide_info.get("title", "")
        
        # 設置副標題
        subtitle = slide.placeholders[1]
        subtitle.text = slide_info.get("subtitle", "")
        
        # 應用樣式
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_toc_slide(self, prs, slide_info):
        """創建目錄幻燈片
        
        Args:
            prs: 演示文稿對象
            slide_info: 幻燈片信息
        """
        # 使用章節幻燈片佈局
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        # 設置標題
        title = slide.shapes.title
        title.text = slide_info.get("title", "目錄")
        
        # 添加目錄項目
        content = slide.placeholders[1]
        tf = content.text_frame
        
        for item in slide_info.get("items", []):
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.space_after = Pt(12)
    
    def _create_content_slide(self, prs, slide_info):
        """創建內容幻燈片
        
        Args:
            prs: 演示文稿對象
            slide_info: 幻燈片信息
        """
        # 使用標題和內容佈局
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 設置標題
        title = slide.shapes.title
        title.text = slide_info.get("title", "")
        
        # 添加要點
        content = slide.placeholders[1]
        tf = content.text_frame
        
        for point in slide_info.get("bullet_points", []):
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 0  # 一級項目符號
    
    def _create_summary_slide(self, prs, slide_info):
        """創建總結幻燈片
        
        Args:
            prs: 演示文稿對象
            slide_info: 幻燈片信息
        """
        # 使用標題和內容佈局
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # 設置標題
        title = slide.shapes.title
        title.text = slide_info.get("title", "總結")
        
        # 設置標題樣式
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        
        # 添加要點
        content = slide.placeholders[1]
        tf = content.text_frame
        
        for point in slide_info.get("bullet_points", []):
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.font.bold = True
    
    def generate_all_presentations(self, translation_dir, output_dir="presentations"):
        """為所有翻譯文件生成演示文稿
        
        Args:
            translation_dir: 翻譯文件目錄
            output_dir: 輸出目錄
            
        Returns:
            生成的演示文稿列表
        """
        os.makedirs(output_dir, exist_ok=True)
        
        generated_files = []
        
        # 查找所有翻譯數據JSON文件
        for filename in os.listdir(translation_dir):
            if filename.endswith("_translation_data.json"):
                json_path = os.path.join(translation_dir, filename)
                
                # 載入翻譯數據
                with open(json_path, 'r', encoding='utf-8') as f:
                    translation_data = json.load(f)
                
                # 創建輸出路徑
                base_name = filename.replace("_translation_data.json", "")
                output_path = os.path.join(output_dir, f"{base_name}_presentation.pptx")
                
                # 生成演示文稿
                logger.info(f"為 {base_name} 生成演示文稿")
                try:
                    self.generate_ppt_from_translation(translation_data, output_path)
                    generated_files.append(output_path)
                except Exception as e:
                    logger.error(f"生成演示文稿時出錯: {str(e)}")
        
        logger.info(f"已生成 {len(generated_files)} 個演示文稿")
        return generated_files

# 使用示例
if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="MCP PPT生成器")
    parser.add_argument("--translation-data", type=str, help="翻譯數據JSON文件路徑")
    parser.add_argument("--output", type=str, default="presentation.pptx", help="輸出PPT路徑")
    parser.add_argument("--translation-dir", type=str, help="翻譯數據目錄，用於批量生成")
    parser.add_argument("--output-dir", type=str, default="presentations", help="批量生成時的輸出目錄")
    args = parser.parse_args()
    
    # 初始化生成器
    generator = MCPPPTGenerator()
    
    if args.translation_data:
        # 載入翻譯數據
        with open(args.translation_data, 'r', encoding='utf-8') as f:
            translation_data = json.load(f)
        
        # 生成演示文稿
        generator.generate_ppt_from_translation(translation_data, args.output)
        print(f"已生成演示文稿：{args.output}")
    elif args.translation_dir:
        # 批量生成
        generated_files = generator.generate_all_presentations(args.translation_dir, args.output_dir)
        print(f"已生成 {len(generated_files)} 個演示文稿")
    else:
        print("請提供翻譯數據文件路徑或翻譯數據目錄")